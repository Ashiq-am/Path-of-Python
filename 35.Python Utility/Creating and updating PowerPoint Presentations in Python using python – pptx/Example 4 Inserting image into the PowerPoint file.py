from pptx import Presentation
from pptx.util import Inches

# Giving Image path
img_path = 'bg_bg.png'

# Creating an Presentation object
ppt = Presentation()

# Selecting blank slide
blank_slide_layout = ppt.slide_layouts[6]

# Attaching slide to ppt
slide = ppt.slides.add_slide(blank_slide_layout)

# For margins
left = top = Inches(1)

# adding images
pic = slide.shapes.add_picture(img_path,
							left, top)

left = Inches(1)
height = Inches(1)

pic = slide.shapes.add_picture(img_path, left,
							top, height = height)
# save file
ppt.save('test_4.pptx')

print("Done")
