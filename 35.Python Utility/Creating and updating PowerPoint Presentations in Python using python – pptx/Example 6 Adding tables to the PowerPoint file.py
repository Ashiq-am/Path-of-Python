# importing
from pptx import Presentation
from pptx.util import Inches

# create a Presentation object
ppt = Presentation()

# Adding a blank slide in out ppt
slide = ppt.slides.add_slide(ppt.slide_layouts[6])

# Adjusting the width !
x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)

# Adding tables
shape = slide.shapes.add_table(3, 4, x,
							y, cx, cy)

# Saving the file
ppt.save("Tabel_Tutorial.pptx")

print("done")
