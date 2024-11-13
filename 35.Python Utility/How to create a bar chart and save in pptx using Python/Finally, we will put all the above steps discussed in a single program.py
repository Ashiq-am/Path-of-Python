import requests
from plotly.graph_objs import Bar
from plotly import offline
from pptx import Presentation
from pptx.util import Inches
from datetime import date


def github_api():
    # siteurl and headers
    site_url = 'https://api.github.com/search/repositories?q=language:javascript&sort=stars'
    headers = {'Accept': 'application/vnd.github.v3+json'}

    # response and parsing the response.
    response = requests.get(site_url, headers=headers)
    response_json = response.json()

    repositories = response_json['items']

    # loop the repositories
    repo_names, repo_stars = [], []
    for repo_info in repositories:
        repo_names.append(repo_info['name'])
        repo_stars.append(repo_info['stargazers_count'])

    # graph plotting
    data_plots = [{'type': 'bar', 'x': repo_names, 'y': repo_stars}]
    layout = {'title': 'GItHubs Most Popular Javascript Projects',
              'xaxis': {'title': 'Repository'},
              'yaxis': {'title': 'Stars'}}

    # saving graph to a Most_Popular_JavaScript_Repos.png
    fig = {'data': data_plots, 'layout': layout}
    offline.plot(fig, image='png', image_filename='Most_Popular_JavaScript_Repos')


def create_pptx_report():
    # create an Object
    ppt = Presentation()
    first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])

    # title (inluded date)
    title = "Popular JavaScript Repositories in GitHub - " + str(date.today())

    # set the title on first slide
    first_slide.shapes[0].text_frame.paragraphs[0].text = title

    # slide 2 - set the image
    img = 'Most_Popular_JavaScript_Repos.png'
    second_slide = ppt.slide_layouts[1]
    slide2 = ppt.slides.add_slide(second_slide)

    # play with the image attributes if you are not OK with the heigth and widthh
    pic = slide2.shapes.add_picture(img, left=Inches(2), top=Inches(1), height=Inches(5))

    # save the powerpoint presentation
    ppt.save('Javascript_report.pptx')


if __name__ == '__main__':
    github_api()
    create_pptx_report()
