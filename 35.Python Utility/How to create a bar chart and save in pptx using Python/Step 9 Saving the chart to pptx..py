from pptx import Presentation
from pptx.util import Inches
from datetime import date

# create an Object
ppt = Presentation()
first_slide = ppt.slides.add_slide(ppt.slide_layouts[0])

# title (inluded date)
title = "Popular JavaScript Repositories in GitHub - " + str(date.today())


# set the title on first slide
first_slide.shapes[0].text_frame.paragraphs[0].text = title

# slide 2 - set the image
img = 'Most_Popular_JavaScript_Repos.png'
second_slide = ppt.slide_layouts[1]
slide2 = ppt.slides.add_slide(second_slide)

# play with the image attributes if you are not OK with the heigth and widthh
pic = slide2.shapes.add_picture(img, left= Inches(2),top = Inches(1),height = Inches(5))

# save the powerpoint presentation
ppt.save('Javascript_report.pptx')
