from pptx import Presentation

# create an object ppt
ppt = Presentation()

# add a new slide
slide = ppt.slides.add_slide(ppt.slide_layouts[0])

# Set the Text to
slide.shapes.title.text = "Popular JavaScript Repositories in GitHub"

# save the powerpoint
ppt.save('Javascript_report.pptx')
